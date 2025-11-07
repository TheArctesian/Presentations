#!/usr/bin/env python3
"""
Nord Dark Theme Presentation Generator
Creates an attractive PowerPoint template with Nord color scheme
Compatible with Google Slides import
"""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Nord color palette
nord = {
    "polar1": RGBColor(46, 52, 64),      # Darkest background
    "polar2": RGBColor(59, 66, 82),      # Dark background
    "polar3": RGBColor(67, 76, 94),      # Medium dark
    "polar4": RGBColor(76, 86, 106),     # Lighter dark
    "snow1": RGBColor(216, 222, 233),    # Body text
    "snow2": RGBColor(229, 233, 240),    # Light text
    "snow3": RGBColor(236, 239, 244),    # Lightest text
    "frost1": RGBColor(143, 188, 187),   # Teal
    "frost2": RGBColor(136, 192, 208),   # Cyan
    "frost3": RGBColor(129, 161, 193),   # Light blue
    "frost4": RGBColor(94, 129, 172),    # Blue
    "aurora_red": RGBColor(191, 97, 106),
    "aurora_orange": RGBColor(208, 135, 112),
    "aurora_yellow": RGBColor(235, 203, 139),
    "aurora_green": RGBColor(163, 190, 140),
    "aurora_purple": RGBColor(180, 142, 173),
}

# Design constants
MARGIN = Inches(0.5)
TITLE_HEIGHT = Inches(1)
ACCENT_LINE_HEIGHT = Inches(0.08)


def add_dark_background(slide, color_key="polar1"):
    """Add a solid dark background to a slide"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = nord[color_key]


def add_accent_line(slide, y_pos, color_key="frost4", width=None):
    """Add a horizontal accent line"""
    if width is None:
        width = Inches(9)
    line = slide.shapes.add_shape(
        1,  # Rectangle
        MARGIN,
        y_pos,
        width,
        ACCENT_LINE_HEIGHT
    )
    line.fill.solid()
    line.fill.fore_color.rgb = nord[color_key]
    line.line.fill.background()
    return line


def add_corner_accent(slide, color_key="frost3"):
    """Add decorative corner element"""
    # Top-right corner shape
    size = Inches(0.3)
    corner = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(9.5),
        Inches(0.2),
        size,
        size
    )
    corner.fill.solid()
    corner.fill.fore_color.rgb = nord[color_key]
    corner.line.fill.background()
    return corner


def style_text_frame(text_frame, text, font_size, color_key, bold=False, alignment=None):
    """Helper to style text frames consistently"""
    text_frame.text = text
    para = text_frame.paragraphs[0]
    para.font.size = Pt(font_size)
    para.font.bold = bold
    para.font.color.rgb = nord[color_key]
    if alignment:
        para.alignment = alignment


def apply_slide_master_backgrounds(prs):
    """Apply Nord dark theme to all slide master layouts"""
    # Apply background to all slide layouts
    for layout in prs.slide_layouts:
        try:
            background = layout.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = nord["polar1"]
        except:
            # Some layouts may not allow background modification
            pass


def create_nord_presentation():
    """Create an attractive Nord-themed presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Try to apply backgrounds to slide master (for Google Slides compatibility)
    apply_slide_master_backgrounds(prs)

    # ============================================================
    # SLIDE 1: Title Slide with Geometric Accents
    # ============================================================
    # Use Title Slide layout (0) for better Google Slides compatibility
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    add_dark_background(slide1)

    # Style the built-in title and subtitle placeholders
    title = slide1.shapes.title
    style_text_frame(title.text_frame, "Nord Dark Theme", 60, "snow3", bold=True, alignment=PP_ALIGN.CENTER)

    # Find subtitle placeholder
    for shape in slide1.placeholders:
        if shape.placeholder_format.idx == 1:  # Subtitle
            style_text_frame(shape.text_frame, "A Modern Presentation Template", 26, "frost2", alignment=PP_ALIGN.CENTER)

    # Large decorative shape on left
    left_accent = slide1.shapes.add_shape(
        1,  # Rectangle
        Inches(0),
        Inches(0),
        Inches(0.3),
        Inches(7.5)
    )
    left_accent.fill.solid()
    left_accent.fill.fore_color.rgb = nord["frost4"]
    left_accent.line.fill.background()

    # Bottom decorative line
    add_accent_line(slide1, Inches(6.8), "frost3")

    # Small corner accents with Aurora colors
    for i, color in enumerate(["aurora_red", "aurora_orange", "aurora_yellow", "aurora_green"]):
        small_box = slide1.shapes.add_shape(
            1,
            Inches(3.5 + i * 0.8),
            Inches(5.2),
            Inches(0.5),
            Inches(0.5)
        )
        small_box.fill.solid()
        small_box.fill.fore_color.rgb = nord[color]
        small_box.line.fill.background()

    # ============================================================
    # SLIDE 2: Section Header
    # ============================================================
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank with title
    add_dark_background(slide2)

    # Bold accent bar on left
    accent_bar = slide2.shapes.add_shape(
        1,
        Inches(0),
        Inches(2.5),
        Inches(0.15),
        Inches(2.5)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = nord["aurora_purple"]
    accent_bar.line.fill.background()

    # Section title
    section_box = slide2.shapes.add_textbox(Inches(0.5), Inches(3), Inches(9), Inches(1.5))
    section_frame = section_box.text_frame
    section_frame.text = "Section Header"
    section_para = section_frame.paragraphs[0]
    section_para.font.size = Pt(54)
    section_para.font.bold = True
    section_para.font.color.rgb = nord["snow3"]

    # Subtitle
    section_subtitle = section_frame.add_paragraph()
    section_subtitle.text = "Use this for major section breaks"
    section_subtitle.font.size = Pt(24)
    section_subtitle.font.color.rgb = nord["frost2"]
    section_subtitle.space_before = Pt(12)

    # ============================================================
    # SLIDE 3: Content Slide with Bullets
    # ============================================================
    # Use Title and Content layout (1) for better compatibility
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    add_dark_background(slide3)

    # Style title
    title3 = slide3.shapes.title
    style_text_frame(title3.text_frame, "Key Features", 44, "snow3", bold=True)

    # Accent line under title
    add_accent_line(slide3, Inches(1.35), "frost4")

    # Corner decoration
    add_corner_accent(slide3, "aurora_green")

    # Style the content placeholder
    content_placeholder = None
    for shape in slide3.placeholders:
        if shape.placeholder_format.type == 2:  # Content/Body
            content_placeholder = shape
            break

    if content_placeholder:
        text_frame = content_placeholder.text_frame
        text_frame.clear()

        bullets = [
            "Clean, minimal design focused on content",
            "Nord color palette for comfortable reading",
            "Perfect for technical presentations",
            "Multiple layout options included",
        ]

        for i, text in enumerate(bullets):
            if i == 0:
                para = text_frame.paragraphs[0]
            else:
                para = text_frame.add_paragraph()
            para.text = text
            para.font.size = Pt(24)
            para.font.color.rgb = nord["snow2"]
            para.space_after = Pt(18)
            para.level = 0

    # Decorative accent box at bottom
    accent_box = slide3.shapes.add_shape(
        1,
        Inches(1),
        Inches(6.2),
        Inches(3),
        Inches(0.8)
    )
    accent_box.fill.solid()
    accent_box.fill.fore_color.rgb = nord["polar3"]
    accent_box.line.color.rgb = nord["frost4"]
    accent_box.line.width = Pt(2)

    note_frame = accent_box.text_frame
    note_frame.text = "Pro Tip: Use Frost colors for accents"
    note_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    note_para = note_frame.paragraphs[0]
    note_para.alignment = PP_ALIGN.CENTER
    note_para.font.size = Pt(18)
    note_para.font.color.rgb = nord["frost2"]

    # ============================================================
    # SLIDE 4: Two-Column Layout
    # ============================================================
    slide4 = prs.slides.add_slide(prs.slide_layouts[3])  # Two content layout
    add_dark_background(slide4)

    # Style title
    title4 = slide4.shapes.title
    style_text_frame(title4.text_frame, "Two-Column Layout", 44, "snow3", bold=True)

    add_accent_line(slide4, Inches(1.35), "frost3")
    add_corner_accent(slide4, "aurora_orange")

    # Left column
    left_col = slide4.shapes.add_textbox(Inches(0.7), Inches(2.2), Inches(4), Inches(4.5))
    left_frame = left_col.text_frame
    left_frame.word_wrap = True

    left_title = left_frame.paragraphs[0]
    left_title.text = "Left Column"
    left_title.font.size = Pt(28)
    left_title.font.bold = True
    left_title.font.color.rgb = nord["frost2"]

    left_content = left_frame.add_paragraph()
    left_content.text = "Use this layout for comparisons, pros/cons, or parallel concepts."
    left_content.font.size = Pt(20)
    left_content.font.color.rgb = nord["snow1"]
    left_content.space_before = Pt(12)

    for item in ["Point one", "Point two", "Point three"]:
        bullet = left_frame.add_paragraph()
        bullet.text = f"• {item}"
        bullet.font.size = Pt(18)
        bullet.font.color.rgb = nord["snow1"]
        bullet.space_before = Pt(8)

    # Vertical divider
    divider = slide4.shapes.add_shape(
        1,
        Inches(5),
        Inches(2),
        Inches(0.03),
        Inches(4.8)
    )
    divider.fill.solid()
    divider.fill.fore_color.rgb = nord["polar4"]
    divider.line.fill.background()

    # Right column
    right_col = slide4.shapes.add_textbox(Inches(5.3), Inches(2.2), Inches(4), Inches(4.5))
    right_frame = right_col.text_frame
    right_frame.word_wrap = True

    right_title = right_frame.paragraphs[0]
    right_title.text = "Right Column"
    right_title.font.size = Pt(28)
    right_title.font.bold = True
    right_title.font.color.rgb = nord["aurora_green"]

    right_content = right_frame.add_paragraph()
    right_content.text = "Each column can have its own accent color for visual distinction."
    right_content.font.size = Pt(20)
    right_content.font.color.rgb = nord["snow1"]
    right_content.space_before = Pt(12)

    for item in ["Feature A", "Feature B", "Feature C"]:
        bullet = right_frame.add_paragraph()
        bullet.text = f"• {item}"
        bullet.font.size = Pt(18)
        bullet.font.color.rgb = nord["snow1"]
        bullet.space_before = Pt(8)

    # ============================================================
    # SLIDE 5: Color Palette Showcase
    # ============================================================
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])  # Title only
    add_dark_background(slide5)

    # Title
    title5 = slide5.shapes.title
    style_text_frame(title5.text_frame, "Nord Color Palette", 44, "snow3", bold=True)

    add_accent_line(slide5, Inches(1.35), "frost2")

    # Color grid
    color_groups = [
        ("Frost", ["frost1", "frost2", "frost3", "frost4"]),
        ("Aurora", ["aurora_red", "aurora_orange", "aurora_yellow", "aurora_green", "aurora_purple"]),
    ]

    y_start = Inches(2.2)

    for group_name, colors in color_groups:
        # Group label
        label_box = slide5.shapes.add_textbox(Inches(0.8), y_start, Inches(2), Inches(0.5))
        label_frame = label_box.text_frame
        label_frame.text = group_name
        label_para = label_frame.paragraphs[0]
        label_para.font.size = Pt(24)
        label_para.font.bold = True
        label_para.font.color.rgb = nord["snow2"]

        # Color swatches
        x_pos = Inches(2.5)
        for color_key in colors:
            swatch = slide5.shapes.add_shape(
                1,
                x_pos,
                y_start,
                Inches(1.2),
                Inches(0.8)
            )
            swatch.fill.solid()
            swatch.fill.fore_color.rgb = nord[color_key]
            swatch.line.fill.background()

            # Color name
            name_box = slide5.shapes.add_textbox(x_pos, y_start + Inches(0.85), Inches(1.2), Inches(0.3))
            name_frame = name_box.text_frame
            name_frame.text = color_key
            name_para = name_frame.paragraphs[0]
            name_para.alignment = PP_ALIGN.CENTER
            name_para.font.size = Pt(10)
            name_para.font.color.rgb = nord["snow1"]

            x_pos += Inches(1.3)

        y_start += Inches(1.5)

    # Usage note at bottom
    note_box = slide5.shapes.add_textbox(Inches(1), Inches(6.3), Inches(8), Inches(0.8))
    note_frame = note_box.text_frame
    note_frame.text = "Use Frost colors for UI elements and links • Use Aurora colors for highlights and callouts"
    note_para = note_frame.paragraphs[0]
    note_para.alignment = PP_ALIGN.CENTER
    note_para.font.size = Pt(16)
    note_para.font.color.rgb = nord["polar4"]

    # ============================================================
    # SLIDE 6: Closing Slide
    # ============================================================
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    add_dark_background(slide6)

    # Decorative top border
    for i in range(5):
        top_box = slide6.shapes.add_shape(
            1,
            Inches(i * 2),
            Inches(0),
            Inches(1.8),
            Inches(0.2)
        )
        top_box.fill.solid()
        colors = ["frost4", "frost3", "frost2", "frost3", "frost4"]
        top_box.fill.fore_color.rgb = nord[colors[i]]
        top_box.line.fill.background()

    # Thank you message
    thank_box = slide6.shapes.add_textbox(Inches(1), Inches(2.8), Inches(8), Inches(1))
    thank_frame = thank_box.text_frame
    thank_frame.text = "Thank You"
    thank_para = thank_frame.paragraphs[0]
    thank_para.alignment = PP_ALIGN.CENTER
    thank_para.font.size = Pt(54)
    thank_para.font.bold = True
    thank_para.font.color.rgb = nord["snow3"]

    # Subtitle
    contact_box = slide6.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(0.6))
    contact_frame = contact_box.text_frame
    contact_frame.text = "Made with Nord Dark Theme"
    contact_para = contact_frame.paragraphs[0]
    contact_para.alignment = PP_ALIGN.CENTER
    contact_para.font.size = Pt(24)
    contact_para.font.color.rgb = nord["frost2"]

    # Bottom decorative elements
    for i, color in enumerate(["aurora_red", "aurora_orange", "aurora_yellow", "aurora_green", "aurora_purple"]):
        circle = slide6.shapes.add_shape(
            1,
            Inches(3 + i * 0.9),
            Inches(5.5),
            Inches(0.6),
            Inches(0.6)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = nord[color]
        circle.line.fill.background()

    # ============================================================
    # Save Presentation
    # ============================================================
    output_path = '/output/nord_dark_theme.pptx'
    prs.save(output_path)

    print("=" * 60)
    print("✓ Nord Dark Theme PowerPoint created successfully!")
    print("=" * 60)
    print(f"📁 Location: {output_path}")
    print(f"📊 Slides: 6")
    print()
    print("Slide Contents:")
    print("  1. Title Slide (with geometric accents)")
    print("  2. Section Header (bold accent bar)")
    print("  3. Content Slide (bullet points + tip box)")
    print("  4. Two-Column Layout (comparison/parallel content)")
    print("  5. Color Palette Showcase (Frost & Aurora)")
    print("  6. Closing Slide (thank you + decorative elements)")
    print()
    print("🎨 Features:")
    print("  • Full Nord color palette integration")
    print("  • Multiple layout options")
    print("  • Decorative accent elements")
    print("  • Professional typography hierarchy")
    print("  • Improved Google Slides compatibility")
    print("=" * 60)


if __name__ == '__main__':
    create_nord_presentation()
